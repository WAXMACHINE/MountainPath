import sys
print(sys.executable)
print(sys.version)


import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
from math import pi
import os
from pptx import Presentation
from pptx.util import Inches
import shutil

from pptx.util import Inches, Pt, Cm
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

def add_chart_to_ppt(template_path, chart_path, output_path, slide_index=3):
    prs = Presentation(template_path)
    slide = prs.slides[slide_index]

    # Set the size of the chart in the PowerPoint slide
    width = Cm(11.65)  # Width
    height = Cm(10.65)  # Height
    left = Cm(12.44)  # Left position in centimeters
    top = Cm(2.22) 

    slide.shapes.add_picture(chart_path, left, top, width, height)
    prs.save(output_path)


import matplotlib.pyplot as plt
import numpy as np
from math import pi

def plot_and_save_radar_chart(name, company_name, individual_scores, average_scores, chart_path, categories):
    N = len(categories)
    angles = [n / float(N) * 2 * pi for n in range(N)]
    angles += angles[:1]
    
    fig, ax = plt.subplots(figsize=(11.74, 10.65), subplot_kw=dict(polar=True))

    # Convert categories to readable names
    readable_categories = [cat.replace('_', ' ').title() for cat in categories]

    # Plot the radar chart
    # Others' average plot (Soft Dark Blue)
    average_scores += average_scores[:1]
    ax.plot(angles, average_scores, linewidth=2, linestyle='solid', color='darkblue', alpha=1)

    # Individual plot (Soft Dark Red)
    individual_scores += individual_scores[:1]
    ax.plot(angles, individual_scores, linewidth=2, linestyle='solid', color='darkred', alpha=1)

    # Set the color and position of the category labels
    ax.set_xticks(angles[:-1])
    ax.set_xticklabels(readable_categories, color='black', size=12)

    # Adjust the plot parameters
    ax.set_rlabel_position(30)  # Move radial labels position


    # Add a title
    plt.title(f"{name} | {company_name}", size=16, color='black', y=1.1)

    plt.savefig(chart_path, bbox_inches='tight')
    plt.close()




# Define the categories
categories = ['apprentissage', 'conscience_de_soi', 'gestion_des_risques', 
              'intelligence_relationnelle', 'mise_en_action', 
              'prise_de_decision', 'prise_de_hauteur']

# Start from the current directory
current_directory = os.getcwd()
company_folders = [f.path for f in os.scandir(current_directory) if f.is_dir()]
template_path = os.path.join(current_directory, 'template.pptx')

for folder in company_folders:
    company_name = os.path.basename(folder)
    individual_path = os.path.join(folder, 'individual.xlsx')
    others_path = os.path.join(folder, 'others.xlsx')
    mapping_path = os.path.join(folder, 'mapping.xlsx')

    individual_df = pd.read_excel(individual_path)
    others_df = pd.read_excel(others_path)
    mapping_df = pd.read_excel(mapping_path)

    for index, row in individual_df.iterrows():
        name = row['Merci de bien vouloir répondre à ce questionnaire individuel. Pour démarrer, pouvez-vous nous indiquer votre prénom et nom ?']
        individual_scores = [row[category] for category in categories]

        # Find the corresponding code for the individual
        code = mapping_df[mapping_df['Who they rate'] == name]['Code'].iloc[0]

        # Filter ratings by others for this individual
        ratings_by_others = others_df[others_df['Code'] == code]
        average_scores = ratings_by_others[categories].mean().tolist()

        # Plot and save radar chart
        chart_path = os.path.join(folder, f"{company_name} - {name}_chart.png")
        plot_and_save_radar_chart(name, company_name, individual_scores, average_scores, chart_path, categories)


        # Add chart to PowerPoint and save
        output_pptx_path = os.path.join(folder, f"{company_name} - {name}.pptx")
        shutil.copy(template_path, output_pptx_path)
        add_chart_to_ppt(output_pptx_path, chart_path, output_pptx_path)

        # Optional: Remove the chart image file if not needed
        os.remove(chart_path)
